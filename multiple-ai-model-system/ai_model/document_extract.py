"""
Utility to extract text from uploaded documents (PDF, DOCX, TXT, etc.)
for AI analysis. Falls back gracefully when libraries are missing.
"""
import os
import io
import base64
import tempfile
import logging

logger = logging.getLogger(__name__)


def extract_text_from_base64(data_uri: str, filename: str = "") -> dict:
    """
    Extract text from a base64-encoded document.
    
    Args:
        data_uri: Either a data URI (data:application/pdf;base64,...) or raw base64 string
        filename: Optional filename hint for format detection
    
    Returns:
        {"text": str, "error": str|None, "pages": int}
    """
    try:
        # Parse data URI
        if data_uri.startswith("data:"):
            header, encoded = data_uri.split(",", 1) if "," in data_uri else ("", data_uri)
            mime_type = header.split(":")[1].split(";")[0] if ":" in header else ""
        else:
            encoded = data_uri
            mime_type = _guess_mime_from_filename(filename)

        # Fix base64 padding
        missing_padding = len(encoded) % 4
        if missing_padding:
            encoded += "=" * (4 - missing_padding)

        raw_bytes = base64.b64decode(encoded)
        return extract_text_from_bytes(raw_bytes, mime_type, filename)

    except Exception as e:
        logger.error(f"Document extraction failed: {e}")
        return {"text": "", "error": str(e), "pages": 0}


def extract_text_from_bytes(raw_bytes: bytes, mime_type: str = "", filename: str = "") -> dict:
    """
    Extract text from raw document bytes.
    """
    if not mime_type:
        mime_type = _guess_mime_from_filename(filename)
    
    mime_lower = mime_type.lower()

    # PDF
    if "pdf" in mime_lower or filename.lower().endswith(".pdf"):
        return _extract_pdf(raw_bytes)

    # DOCX
    if "wordprocessingml" in mime_lower or "msword" in mime_lower or \
       "vnd.openxmlformats" in mime_lower or filename.lower().endswith(".docx"):
        return _extract_docx(raw_bytes)

    # Plain text variants
    if any(x in mime_lower for x in ["text/", "json", "csv", "xml", "html"]) or \
       filename.lower().endswith((".txt", ".md", ".csv", ".json", ".xml", ".html",
                                   ".log", ".py", ".js", ".ts", ".java", ".sql",
                                   ".yml", ".yaml", ".tsx", ".jsx")):
        return _extract_plain_text(raw_bytes)

    # PPT/PPTX - basic support
    if "presentation" in mime_lower or filename.lower().endswith((".pptx", ".ppt")):
        return _extract_pptx(raw_bytes)

    # XLSX - basic support
    if "spreadsheet" in mime_lower or filename.lower().endswith((".xlsx", ".xls")):
        return _extract_xlsx(raw_bytes)

    return {"text": "", "error": f"Unsupported document format: {mime_type or filename}", "pages": 0}


def _extract_pdf(raw_bytes: bytes) -> dict:
    """Extract text from PDF using PyPDF2 or pdfplumber."""
    # Try PyPDF2 first
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(raw_bytes))
        pages = len(reader.pages)
        text_parts = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"[Page {i+1}]\n{page_text.strip()}")
        
        full_text = "\n\n".join(text_parts)
        if full_text.strip():
            return {"text": full_text, "error": None, "pages": pages}
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"PyPDF2 extraction failed: {e}")

    # Fallback: pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
            pages = len(pdf.pages)
            text_parts = []
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[Page {i+1}]\n{page_text.strip()}")
            
            full_text = "\n\n".join(text_parts)
            if full_text.strip():
                return {"text": full_text, "error": None, "pages": pages}
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"pdfplumber extraction failed: {e}")

    return {
        "text": "",
        "error": "PDF extraction requires PyPDF2 or pdfplumber. Install with: pip install PyPDF2 pdfplumber",
        "pages": 0
    }


def _extract_docx(raw_bytes: bytes) -> dict:
    """Extract text from DOCX using python-docx."""
    try:
        import docx
        doc = docx.Document(io.BytesIO(raw_bytes))
        text_parts = [para.text for para in doc.paragraphs if para.text.strip()]

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        full_text = "\n".join(text_parts)
        return {"text": full_text, "error": None, "pages": 1}
    except ImportError:
        return {
            "text": "",
            "error": "DOCX extraction requires python-docx. Install with: pip install python-docx",
            "pages": 0
        }
    except Exception as e:
        return {"text": "", "error": f"DOCX extraction failed: {str(e)}", "pages": 0}


def _extract_plain_text(raw_bytes: bytes) -> dict:
    """Extract text from plain text files."""
    for encoding in ["utf-8", "latin-1", "cp1252"]:
        try:
            text = raw_bytes.decode(encoding)
            return {"text": text.strip(), "error": None, "pages": 1}
        except (UnicodeDecodeError, ValueError):
            continue
    return {"text": "", "error": "Could not decode text file", "pages": 0}


def _extract_pptx(raw_bytes: bytes) -> dict:
    """Extract text from PPTX."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(raw_bytes))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
            if slide_texts:
                text_parts.append(f"[Slide {i+1}]\n" + "\n".join(slide_texts))
        
        return {"text": "\n\n".join(text_parts), "error": None, "pages": len(prs.slides)}
    except ImportError:
        return {"text": "", "error": "PPTX extraction requires python-pptx. Install with: pip install python-pptx", "pages": 0}
    except Exception as e:
        return {"text": "", "error": f"PPTX extraction failed: {str(e)}", "pages": 0}


def _extract_xlsx(raw_bytes: bytes) -> dict:
    """Extract text from XLSX."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return {"text": "\n\n".join(text_parts), "error": None, "pages": len(wb.sheetnames)}
    except ImportError:
        return {"text": "", "error": "XLSX extraction requires openpyxl. Install with: pip install openpyxl", "pages": 0}
    except Exception as e:
        return {"text": "", "error": f"XLSX extraction failed: {str(e)}", "pages": 0}


def _guess_mime_from_filename(filename: str) -> str:
    """Guess MIME type from filename extension."""
    if not filename:
        return ""
    ext = os.path.splitext(filename)[1].lower()
    mime_map = {
        ".pdf": "application/pdf",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".txt": "text/plain",
        ".csv": "text/csv",
        ".json": "application/json",
        ".xml": "text/xml",
        ".html": "text/html",
        ".md": "text/markdown",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
    return mime_map.get(ext, "")
